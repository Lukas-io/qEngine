import os
import pytesseract
import pdfplumber
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from PIL import Image
import numpy as np
import cv2
import docx
import pptx

def get_file_type(file_path):
    """Automatically detects file type based on extension."""
    return os.path.splitext(file_path)[1].lower()

### PDF HANDLING ###
def extract_text_from_pdf(pdf_path, password=None, max_pages=None):
    """Extracts text from a PDF, handling OCR if necessary."""
    with pdfplumber.open(pdf_path) as pdf:
        if password:
            pdf.decrypt(password)

        text = ""
        for i, page in enumerate(pdf.pages):
            if max_pages and i >= max_pages:
                break
            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"  # Extract text-based PDFs
            else:
                image = page.to_image().annotated  # Convert page to an image
                text += pytesseract.image_to_string(image) + "\n"  # OCR for images

    return text.strip()

### IMAGE HANDLING ###
def preprocess_image(image):
    """Enhances an image for better OCR accuracy."""
    image = image.convert("L")  # Convert to grayscale
    img_np = np.array(image)

    # Apply adaptive thresholding
    img_np = cv2.adaptiveThreshold(img_np, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                   cv2.THRESH_BINARY, 11, 2)
    
    return Image.fromarray(img_np)

def extract_text_from_image(image_path):
    """Extracts text from an image."""
    image = Image.open(image_path)
    image = preprocess_image(image)
    return pytesseract.image_to_string(image).strip()

### WORD DOCUMENT HANDLING ###
def extract_text_from_docx(docx_path):
    """Extracts text from a Word document."""
    doc = docx.Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

### POWERPOINT HANDLING ###
def extract_text_from_pptx(pptx_path):
    """Extracts text from a PowerPoint presentation."""
    presentation = pptx.Presentation(pptx_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

### TEXT FILE HANDLING ###
def extract_text_from_txt(txt_path):
    """Extracts text from a plain text file."""
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read().strip()

### UNIVERSAL FUNCTION ###
def extract_text(file_path):
    """Automatically detects the file type and extracts text accordingly."""
    file_type = get_file_type(file_path)

    if file_type == ".pdf":
        return extract_text_from_pdf(file_path)
    elif file_type in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(file_path)
    elif file_type == ".docx":
        return extract_text_from_docx(file_path)
    elif file_type == ".pptx":
        return extract_text_from_pptx(file_path)
    elif file_type == ".txt":
        return extract_text_from_txt(file_path)
    
    else:
        raise ValueError(f"Unsupported file format: {file_type}")